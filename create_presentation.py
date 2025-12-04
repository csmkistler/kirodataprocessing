"""
Professional PowerPoint Presentation Generator
Creates a customer presentation for Signal Processing Visualization Platform
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

def create_presentation():
    # Create presentation object
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Define color scheme (Professional blue theme)
    PRIMARY_COLOR = RGBColor(0, 102, 204)      # Blue
    SECONDARY_COLOR = RGBColor(51, 51, 51)     # Dark gray
    ACCENT_COLOR = RGBColor(0, 176, 80)        # Green
    LIGHT_BG = RGBColor(240, 248, 255)         # Light blue
    WHITE = RGBColor(255, 255, 255)
    
    def add_title_slide(title, subtitle):
        """Add a title slide with gradient background"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        
        # Add background rectangle
        bg = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = PRIMARY_COLOR
        bg.line.fill.background()
        
        # Add title
        title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(54)
        title_para.font.bold = True
        title_para.font.color.rgb = WHITE
        title_para.alignment = PP_ALIGN.CENTER
        
        # Add subtitle
        subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(8), Inches(1))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = subtitle
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.font.size = Pt(28)
        subtitle_para.font.color.rgb = WHITE
        subtitle_para.alignment = PP_ALIGN.CENTER
        
        return slide
    
    def add_content_slide(title, content_items, layout_type="bullet"):
        """Add a content slide with title and bullet points or two columns"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        
        # Add title bar
        title_bar = slide.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(1))
        title_bar.fill.solid()
        title_bar.fill.fore_color.rgb = PRIMARY_COLOR
        title_bar.line.fill.background()
        
        # Add title text
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.6))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(36)
        title_para.font.bold = True
        title_para.font.color.rgb = WHITE
        
        # Add content based on layout type
        if layout_type == "bullet":
            content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(5.5))
            text_frame = content_box.text_frame
            text_frame.word_wrap = True
            
            for i, item in enumerate(content_items):
                if i > 0:
                    text_frame.add_paragraph()
                p = text_frame.paragraphs[i]
                p.text = item
                p.level = 0
                p.font.size = Pt(20)
                p.font.color.rgb = SECONDARY_COLOR
                p.space_before = Pt(12)
        
        elif layout_type == "two_column":
            # Left column
            left_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(4), Inches(5.5))
            left_frame = left_box.text_frame
            left_frame.word_wrap = True
            
            for i, item in enumerate(content_items[:len(content_items)//2]):
                if i > 0:
                    left_frame.add_paragraph()
                p = left_frame.paragraphs[i]
                p.text = item
                p.font.size = Pt(18)
                p.font.color.rgb = SECONDARY_COLOR
                p.space_before = Pt(10)
            
            # Right column
            right_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.5), Inches(4), Inches(5.5))
            right_frame = right_box.text_frame
            right_frame.word_wrap = True
            
            for i, item in enumerate(content_items[len(content_items)//2:]):
                if i > 0:
                    right_frame.add_paragraph()
                p = right_frame.paragraphs[i]
                p.text = item
                p.font.size = Pt(18)
                p.font.color.rgb = SECONDARY_COLOR
                p.space_before = Pt(10)
        
        return slide
    
    def add_feature_slide(title, features):
        """Add a slide with feature boxes"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Add title bar
        title_bar = slide.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(1))
        title_bar.fill.solid()
        title_bar.fill.fore_color.rgb = PRIMARY_COLOR
        title_bar.line.fill.background()
        
        # Add title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.6))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(36)
        title_para.font.bold = True
        title_para.font.color.rgb = WHITE
        
        # Add feature boxes (2x2 grid)
        box_width = Inches(4)
        box_height = Inches(2.5)
        margin = Inches(0.5)
        
        positions = [
            (Inches(0.8), Inches(1.5)),
            (Inches(5.2), Inches(1.5)),
            (Inches(0.8), Inches(4.2)),
            (Inches(5.2), Inches(4.2))
        ]
        
        for i, (feature_title, feature_desc) in enumerate(features[:4]):
            x, y = positions[i]
            
            # Feature box background
            box = slide.shapes.add_shape(1, x, y, box_width, box_height)
            box.fill.solid()
            box.fill.fore_color.rgb = LIGHT_BG
            box.line.color.rgb = PRIMARY_COLOR
            box.line.width = Pt(2)
            
            # Feature title
            title_box = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.2), 
                                                  box_width - Inches(0.4), Inches(0.5))
            title_frame = title_box.text_frame
            title_frame.text = feature_title
            title_para = title_frame.paragraphs[0]
            title_para.font.size = Pt(22)
            title_para.font.bold = True
            title_para.font.color.rgb = PRIMARY_COLOR
            
            # Feature description
            desc_box = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.8), 
                                                 box_width - Inches(0.4), box_height - Inches(1))
            desc_frame = desc_box.text_frame
            desc_frame.text = feature_desc
            desc_frame.word_wrap = True
            desc_para = desc_frame.paragraphs[0]
            desc_para.font.size = Pt(16)
            desc_para.font.color.rgb = SECONDARY_COLOR
        
        return slide
    
    # Slide 1: Title Slide
    add_title_slide(
        "Signal Processing Visualization Platform",
        "Making Complex Signal Analysis Simple, Visual, and Accessible"
    )
    
    # Slide 2: The Problem
    add_content_slide(
        "The Challenge",
        [
            "❌ Complex tools with steep learning curves",
            "❌ No real-time visualization of processing effects",
            "❌ Difficult to experiment with different parameters",
            "❌ Poor integration between generation, processing, and analysis",
            "",
            "✅ Our Solution: An all-in-one platform that makes signal processing",
            "    visual, interactive, and accessible"
        ]
    )
    
    # Slide 3: Key Features Overview
    add_feature_slide(
        "Key Features",
        [
            ("🎵 Multi-Signal Generation", "Generate sine, square, sawtooth waves, and white noise with configurable parameters"),
            ("⚡ Real-Time Processing", "Apply filters and transformations with instant visual feedback"),
            ("📊 Interactive Visualization", "Professional charting with zoom, pan, and side-by-side comparison"),
            ("🔔 Smart Event Detection", "Monitor signals with configurable thresholds and automatic event logging")
        ]
    )
    
    # Slide 4: Signal Generation
    add_content_slide(
        "Multi-Signal Generation",
        [
            "🎵 Sine Waves - Pure tones for testing",
            "⬛ Square Waves - Digital signal simulation",
            "📐 Sawtooth Waves - Audio synthesis",
            "📡 White Noise - System testing",
            "",
            "Configure:",
            "  • Frequency, amplitude, phase",
            "  • Duration and sample rate",
            "  • Instant parameter validation"
        ]
    )
    
    # Slide 5: Signal Processing
    add_content_slide(
        "Real-Time Signal Processing",
        [
            "🔽 Low-Pass Filters",
            "  Remove high-frequency noise",
            "",
            "🔼 High-Pass Filters",
            "  Eliminate DC offset and low-frequency drift",
            "",
            "🎚️ Band-Pass Filters",
            "  Isolate specific frequency ranges",
            "",
            "📈 Gain Adjustment",
            "  Amplify or attenuate signals",
            "",
            "✨ See effects immediately with side-by-side comparison"
        ]
    )
    
    # Slide 6: Technical Excellence
    add_content_slide(
        "Enterprise-Grade Architecture",
        [
            "🏗️ Onion Architecture",
            "  • Core Domain Layer - Pure business logic",
            "  • Application Layer - Use case orchestration",
            "  • Infrastructure Layer - Database & services",
            "  • Presentation Layer - REST API + React UI",
            "",
            "💻 Modern Technology Stack",
            "  • Backend: .NET 10 with ASP.NET Core",
            "  • Frontend: React 18 + TypeScript",
            "  • Time-Series DB: InfluxDB (10:1 compression)",
            "  • Metadata DB: MongoDB (flexible storage)",
            "  • API Documentation: Swagger/OpenAPI"
        ]
    )
    
    # Slide 7: Live Demo Flow
    add_content_slide(
        "Live Demo Flow (3 Minutes)",
        [
            "1️⃣ Generate a 1kHz Sine Wave (30 sec)",
            "   • Show parameter panel",
            "   • Click generate",
            "   • Watch waveform appear instantly",
            "",
            "2️⃣ Apply Low-Pass Filter at 500Hz (30 sec)",
            "   • Configure filter parameters",
            "   • Apply processing",
            "   • See both signals side-by-side",
            "",
            "3️⃣ Test Threshold Detection (20 sec)",
            "   • Set threshold to 0.5",
            "   • Input values and trigger events",
            "",
            "4️⃣ Demonstrate Persistence (20 sec)",
            "   • Show historical signal selector"
        ]
    )
    
    # Slide 8: Use Cases
    add_feature_slide(
        "Use Cases Across Industries",
        [
            ("🔬 Research & Development", "Test algorithms, validate filter designs, prototype audio/RF systems"),
            ("🎓 Education", "Teach signal processing concepts with interactive visual demonstrations"),
            ("🏭 Quality Control", "Monitor production signals, detect anomalies, analyze historical data"),
            ("🎵 Audio Engineering", "Synthesize test tones, analyze frequency response, validate filter designs")
        ]
    )
    
    # Slide 9: Competitive Advantage
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title bar
    title_bar = slide.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(1))
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = PRIMARY_COLOR
    title_bar.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = "What Makes Us Different"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = WHITE
    
    # Comparison table
    table_data = [
        ["Feature", "Traditional Tools", "Our Platform"],
        ["Learning Curve", "Weeks", "✅ Minutes"],
        ["Real-Time Feedback", "❌ No", "✅ Instant"],
        ["Visual Comparison", "Manual", "✅ Side-by-side"],
        ["Setup Complexity", "High", "✅ One-click"],
        ["Cost", "$$$$", "✅ Local & Free"],
        ["Extensibility", "Limited", "✅ Open Architecture"]
    ]
    
    table = slide.shapes.add_table(len(table_data), 3, Inches(1), Inches(1.8), 
                                    Inches(8), Inches(4.5)).table
    
    # Set column widths
    table.columns[0].width = Inches(2.5)
    table.columns[1].width = Inches(2.75)
    table.columns[2].width = Inches(2.75)
    
    # Fill table
    for i, row_data in enumerate(table_data):
        for j, cell_text in enumerate(row_data):
            cell = table.cell(i, j)
            cell.text = cell_text
            
            # Header row styling
            if i == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = PRIMARY_COLOR
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.size = Pt(18)
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = WHITE
                    paragraph.alignment = PP_ALIGN.CENTER
            else:
                # Data rows
                if j == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = LIGHT_BG
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.size = Pt(16)
                    paragraph.font.color.rgb = SECONDARY_COLOR
                    paragraph.alignment = PP_ALIGN.CENTER
    
    # Slide 10: Business Value
    add_content_slide(
        "Business Value",
        [
            "For Organizations:",
            "  📈 Reduce development time by 60%",
            "  💰 Lower training costs with intuitive UI",
            "  🤝 Better collaboration with visual results",
            "  ✅ Built-in quality assurance",
            "",
            "For Individuals:",
            "  🚀 Learn faster with visual feedback",
            "  🔬 Experiment freely with instant results",
            "  🏆 Professional-grade algorithms",
            "  🔒 Complete control - all data stays local"
        ]
    )
    
    # Slide 11: Future Roadmap
    add_content_slide(
        "Future Roadmap",
        [
            "Phase 2 Enhancements:",
            "",
            "📊 Frequency-domain visualization (FFT spectrum)",
            "💾 Export signals to CSV/JSON",
            "📥 Import real-world signal data",
            "🎛️ Additional filter types (notch, all-pass)",
            "🎤 Real-time audio input streaming",
            "➕ Multi-signal arithmetic operations",
            "",
            "Your feedback shapes our roadmap!"
        ]
    )
    
    # Slide 12: Quick Stats
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title bar
    title_bar = slide.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(1))
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = PRIMARY_COLOR
    title_bar.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = "Quick Stats"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = WHITE
    
    # Stats boxes
    stats = [
        ("⚡ <100ms", "Visualization Updates"),
        ("🗄️ 10:1", "Data Compression"),
        ("🎨 4 Types", "Signal Generation"),
        ("🔧 4 Operations", "Signal Processing"),
        ("🏗️ 4 Layers", "Clean Architecture"),
        ("✅ 100%", "Local Deployment")
    ]
    
    positions = [
        (Inches(0.8), Inches(1.8)),
        (Inches(3.7), Inches(1.8)),
        (Inches(6.6), Inches(1.8)),
        (Inches(0.8), Inches(4.3)),
        (Inches(3.7), Inches(4.3)),
        (Inches(6.6), Inches(4.3))
    ]
    
    for (stat, label), (x, y) in zip(stats, positions):
        # Stat box
        box = slide.shapes.add_shape(1, x, y, Inches(2.5), Inches(2))
        box.fill.solid()
        box.fill.fore_color.rgb = ACCENT_COLOR
        box.line.fill.background()
        
        # Stat number
        stat_box = slide.shapes.add_textbox(x, y + Inches(0.3), Inches(2.5), Inches(0.8))
        stat_frame = stat_box.text_frame
        stat_frame.text = stat
        stat_para = stat_frame.paragraphs[0]
        stat_para.font.size = Pt(32)
        stat_para.font.bold = True
        stat_para.font.color.rgb = WHITE
        stat_para.alignment = PP_ALIGN.CENTER
        
        # Label
        label_box = slide.shapes.add_textbox(x, y + Inches(1.2), Inches(2.5), Inches(0.6))
        label_frame = label_box.text_frame
        label_frame.text = label
        label_frame.word_wrap = True
        label_para = label_frame.paragraphs[0]
        label_para.font.size = Pt(16)
        label_para.font.color.rgb = WHITE
        label_para.alignment = PP_ALIGN.CENTER
    
    # Slide 13: Call to Action
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Background
    bg = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = PRIMARY_COLOR
    bg.line.fill.background()
    
    # Main heading
    heading_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(1))
    heading_frame = heading_box.text_frame
    heading_frame.text = "Ready to Transform Your Workflow?"
    heading_para = heading_frame.paragraphs[0]
    heading_para.font.size = Pt(44)
    heading_para.font.bold = True
    heading_para.font.color.rgb = WHITE
    heading_para.alignment = PP_ALIGN.CENTER
    
    # Steps
    steps_box = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(6), Inches(2.5))
    steps_frame = steps_box.text_frame
    steps_frame.word_wrap = True
    
    steps_text = [
        "1. Clone the repository",
        "2. Run start.ps1 (Windows)",
        "3. Start processing in under 60 seconds"
    ]
    
    for i, step in enumerate(steps_text):
        if i > 0:
            steps_frame.add_paragraph()
        p = steps_frame.paragraphs[i]
        p.text = step
        p.font.size = Pt(28)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        p.space_before = Pt(20)
    
    # Contact info
    contact_box = slide.shapes.add_textbox(Inches(2), Inches(6), Inches(6), Inches(0.8))
    contact_frame = contact_box.text_frame
    contact_frame.text = "Let's discuss how this platform can solve\nyour signal processing challenges"
    contact_para = contact_frame.paragraphs[0]
    contact_para.font.size = Pt(20)
    contact_para.font.color.rgb = WHITE
    contact_para.alignment = PP_ALIGN.CENTER
    
    # Slide 14: Thank You
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Background
    bg = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = ACCENT_COLOR
    bg.line.fill.background()
    
    # Thank you text
    thanks_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(2))
    thanks_frame = thanks_box.text_frame
    thanks_frame.text = "Thank You!"
    thanks_para = thanks_frame.paragraphs[0]
    thanks_para.font.size = Pt(72)
    thanks_para.font.bold = True
    thanks_para.font.color.rgb = WHITE
    thanks_para.alignment = PP_ALIGN.CENTER
    
    # Tagline
    tagline_box = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(1))
    tagline_frame = tagline_box.text_frame
    tagline_frame.text = "Signal Processing Visualization Platform"
    tagline_para = tagline_frame.paragraphs[0]
    tagline_para.font.size = Pt(32)
    tagline_para.font.color.rgb = WHITE
    tagline_para.alignment = PP_ALIGN.CENTER
    
    # Save presentation
    prs.save('Signal_Processing_Platform_Presentation.pptx')
    print("✅ Presentation created successfully: Signal_Processing_Platform_Presentation.pptx")
    print("📊 Total slides: 14")
    print("🎨 Professional blue theme applied")
    print("⏱️ Optimized for 3-minute delivery")

if __name__ == "__main__":
    create_presentation()
